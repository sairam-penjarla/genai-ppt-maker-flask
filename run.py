from src.utils import Utilities
from custom_logger import logger
import os
from flask import jsonify
from pptx import Presentation
import json
from flask import (
    Flask, 
    jsonify, 
    render_template, 
    request
)
from src.prompt_templates import ( 
    SLIDE_PLANNING_GUIDELINES,
    SLIDE_CONTENT_GUIDELINES
)
import traceback

app = Flask(__name__)
utils = Utilities()


@app.route("/")
def landing_page():
    previous_session_meta_data = utils.session_utils.get_session_meta_data()
    return render_template("index.html", previous_session_meta_data = previous_session_meta_data)


@app.route('/get_session_data', methods=['POST'])
def get_session_data():
    data = request.get_json()
    session_id = data.get('sessionId')

    session_data = utils.session_utils.get_session_data(session_id)

    return jsonify({'session_data': session_data})


@app.route('/generate_slides', methods=['POST'])
def generate_slides():
    try:
        data = request.get_json()
        user_input = data.get('user_input')
        session_id = data.get('session_id')
        
        messages = utils.get_previous_messages(session_id, SLIDE_PLANNING_GUIDELINES, "slides_planning")
        messages.append({
            "role": "user",
            "content": user_input
        })
        logger.info(messages)
        
        # slide_planning = messages[-2]['content'] # testing
        slide_planning = utils.invoke_llm(messages)
        session_icon = utils.get_session_icon(session_id)
        logger.info(slide_planning)
        slides_content = []
        for slide_details in json.loads(slide_planning)['slide_details']:
            # Fetching the previous messages for slide content
            messages = utils.get_previous_messages(session_id, SLIDE_CONTENT_GUIDELINES, "slides_content")
            messages.append({
                "role": "user",
                "content": slide_details['purpose']
            })
            
            # single_slide_content = messages[-2]['content'] # testing
            single_slide_content = utils.invoke_llm(messages)

            messages.append({
                "role": "assistant",
                "content": single_slide_content
            })
            
            slides_content.append(json.loads(single_slide_content))
        
        return jsonify({'slides_content': slides_content, 'slide_planning': slide_planning, 'session_icon' : session_icon}), 200

    except Exception as e:
        # Print the exception traceback to the console
        print("Error occurred:", e)
        traceback.print_exc()
        return jsonify({"error": "An error occurred while generating slides."}), 500


@app.route('/delete_session', methods=['POST'])
def delete_session():
    data = request.get_json()
    session_id = data.get('session_id')

    if not session_id:
        return jsonify({"error": "Session ID is required."}), 400

    response, status_code = utils.session_utils.delete_session(session_id)
    return jsonify(response), status_code


@app.route('/delete_all_sessions', methods=['POST'])
def delete_all_sessions():
    response, status_code = utils.session_utils.delete_all_sessions()
    return jsonify(response), status_code


def update_single_slide_ppt(thumbnail_id, slide_data, output_folder="static/output"):
    """
    Updates a single-slide PPT file with new slide data and saves the updated file.
    
    Args:
        thumbnail_id (str): The name of the PPT file (without extension) located in static/files.
        slide_data (dict): The slide data to update (contains a single slide's content).
        output_folder (str): The folder where the updated PPT file should be saved.
    
    Returns:
        dict: A dictionary containing the file path of the updated PPT.
    """
    try:
        # Paths
        input_ppt_path = f"static/files/{thumbnail_id}.pptx"
        if not os.path.exists(input_ppt_path):
            return jsonify({"error": "Template PPT file not found."}), 404
        if not os.path.exists(output_folder):
            os.makedirs(output_folder)
        # Output PPT path
        output_ppt_path = os.path.join(output_folder, f"{thumbnail_id}_updated.pptx")
        # Load the PPT file
        presentation = Presentation(input_ppt_path)
        # Ensure it only has one slide
        if len(presentation.slides) != 1:
            return jsonify({"error": "The template PPT file must contain exactly one slide."}), 400
        # Access the first slide
        slide = presentation.slides[0]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            shape_name = shape.name
            if shape_name in slide_data.keys():
                data = slide_data[shape_name]
                text_frame = shape.text_frame
                # Update text content in the shape
                for paragraph in text_frame.paragraphs:
                    for i, run in enumerate(paragraph.runs):
                        if i == 0:
                            run.text = data  # Update the text
                        else:
                            run.text = ''  # Clear other runs
        # Save the updated PPT
        presentation.save(output_ppt_path)

        # Return the path to the updated file
        return jsonify({"file_path": output_ppt_path}), 200

    except Exception as e:
        # Handle errors
        print("Error updating PPT:", str(e))
        return jsonify({"error": "An error occurred while updating the PPT.", "details": str(e)}), 500


def transform_slide_data(input_data):
    input_data = json.loads(input_data)
    # Extract titles and content sections
    titles = input_data.get("titles", {})
    content = input_data.get("content", [])

    # Initialize the transformed data with titles
    transformed_data = {
        "slide_title": titles.get("slide_title", ""),
        "long_title": titles.get("long_title", ""),
        "sub_title": titles.get("sub_title", ""),
    }

    # Add content sections as heading and content pairs
    for i, section in enumerate(content, start=1):
        heading_key = f"heading{i}"
        content_key = f"content{i}"
        transformed_data[heading_key] = section.get("heading", "")
        transformed_data[content_key] = section.get("content", "")

    return transformed_data

@app.route('/get_thumbnails', methods=['POST'])
def get_thumbnails():
    files = os.listdir('static/thumbnails')
    return jsonify({"files": files}), 200

@app.route('/download_slide', methods=['POST'])
def download_slide():
    try:
        data = request.get_json()
        thumbnail_id = data.get('thumbnail_id').replace(".png", "")
        slide_data = data.get('slideData')

        if not thumbnail_id or not slide_data:
            return jsonify({"error": "Invalid input. 'thumbnail_id' and 'slideData' are required."}), 400
        return update_single_slide_ppt(thumbnail_id, transform_slide_data(slide_data))
    except Exception as e:
        # Print the exception traceback to the console
        print("Error occurred:", e)
        traceback.print_exc()
        return jsonify({"error": "An error occurred while generating slides."}), 500


@app.route('/update_session', methods=['POST'])
def update_session():

    data = request.get_json()

    session_id      = data.get('session_id')
    prompt          = data.get('prompt')
    slide_planning  = data.get('slide_planning')
    session_icon  = data.get('session_icon')
    slides_content  = json.dumps(data.get('slides_content'))

    utils.session_utils.add_slides(session_id, prompt, slide_planning, slides_content, session_icon)

    return jsonify({"success" : True})

if __name__ == "__main__":

    app.run(host="0.0.0.0", port=8000)